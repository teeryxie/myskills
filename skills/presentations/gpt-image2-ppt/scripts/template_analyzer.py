#!/usr/bin/env python3
"""
PPT 模板剖析器 -- 把用户的 .pptx + 渲染图喂给 vision，吐出 TemplateProfile JSON。

主要职责：
1. read_pptx_metadata: 用 python-pptx 抽 theme/colors/fonts（可选 fallback）
2. load_template_images: 从目录加载每页 PNG（按文件名排序）
3. compute_source_hash: 算 .pptx + images 的内容哈希做缓存 key
4. VisionClient: 封装 OpenAI 兼容 chat completions 多模态调用
5. vision_analyze: 调 vision 抽 global_style + 每页 summary/page_type/json_schema
6. analyze_template: 顶层入口，带缓存
7. assign_layouts: 按 layout_id / page_type 分配 slide -> layout，优先不复用 layout
8. coerce_fields: 把自由 content 字符串拆解成 layout schema 的 fields
"""
from __future__ import annotations

import base64
import hashlib
import json
import os
import re
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import requests


VISION_TIMEOUT_SECS = int(os.getenv("VISION_TIMEOUT_SECS", "90"))
VISION_MAX_RETRIES = int(os.getenv("VISION_MAX_RETRIES", "2"))
VISION_RETRY_DELAY = int(os.getenv("VISION_RETRY_DELAY", "4"))

CACHE_DIR_NAME = "template_cache"
PROFILE_VERSION = "1"

VALID_PAGE_TYPES = {"cover", "agenda", "section", "content", "data", "quote", "closing", "other"}


# =============================================================================
# 工具
# =============================================================================

def _read_image_b64(path: str) -> str:
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode("ascii")


def _data_url_for(path: str) -> str:
    suffix = Path(path).suffix.lower().lstrip(".") or "png"
    if suffix == "jpg":
        suffix = "jpeg"
    return f"data:image/{suffix};base64,{_read_image_b64(path)}"


def _strip_json_fence(text: str) -> str:
    s = text.strip()
    if s.startswith("```"):
        # 去掉 ```lang ... ```
        s = re.sub(r"^```[a-zA-Z0-9_-]*\s*", "", s)
        if s.endswith("```"):
            s = s[: -3]
        s = s.strip()
    return s


def _parse_json_loose(text: str) -> Any:
    """尽量从 LLM 文本里拽出 JSON。"""
    s = _strip_json_fence(text)
    try:
        return json.loads(s)
    except Exception:
        pass
    # 找第一个 { 和最后一个 } 之间的内容
    start = s.find("{")
    end = s.rfind("}")
    if start != -1 and end != -1 and end > start:
        try:
            return json.loads(s[start : end + 1])
        except Exception:
            pass
    # 找第一个 [ 和最后一个 ]
    start = s.find("[")
    end = s.rfind("]")
    if start != -1 and end != -1 and end > start:
        try:
            return json.loads(s[start : end + 1])
        except Exception:
            pass
    raise ValueError(f"无法解析为 JSON: {text[:300]}")


# =============================================================================
# .pptx 元数据
# =============================================================================

def read_pptx_metadata(pptx_path: str) -> Dict[str, Any]:
    """用 python-pptx 抽 theme/layouts/字体/配色。失败返回空 dict。"""
    try:
        from pptx import Presentation
    except ImportError:
        print("(!)  python-pptx 未安装，跳过 .pptx 结构读取")
        return {}

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"(!)  打开 .pptx 失败: {e}")
        return {}

    width_emu = prs.slide_width
    height_emu = prs.slide_height
    aspect = "16:9"
    if width_emu and height_emu:
        ratio = width_emu / height_emu
        if abs(ratio - 16 / 9) < 0.05:
            aspect = "16:9"
        elif abs(ratio - 4 / 3) < 0.05:
            aspect = "4:3"
        else:
            aspect = f"{width_emu}:{height_emu}"

    layouts = []
    try:
        for layout in prs.slide_layouts:
            layouts.append({"name": getattr(layout, "name", "") or "", "placeholder_count": len(layout.placeholders)})
    except Exception:
        pass

    page_count = 0
    try:
        page_count = len(prs.slides)
    except Exception:
        pass

    return {
        "aspect": aspect,
        "slide_width_emu": width_emu,
        "slide_height_emu": height_emu,
        "page_count": page_count,
        "master_layouts": layouts,
    }


# =============================================================================
# 图片加载 + 哈希
# =============================================================================

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}


def load_template_images(images_dir: str) -> List[str]:
    """按文件名升序加载 images_dir 下所有图片，返回绝对路径列表。"""
    if not images_dir:
        return []
    p = Path(images_dir)
    if not p.is_dir():
        raise FileNotFoundError(f"模板图片目录不存在: {images_dir}")
    files = sorted(
        [str(f) for f in p.iterdir() if f.suffix.lower() in IMAGE_EXTS],
        key=lambda s: s.lower(),
    )
    return files


def compute_source_hash(pptx_path: Optional[str], images: List[str]) -> str:
    """对 .pptx + 所有图片内容计算 sha256，作为缓存 key。"""
    h = hashlib.sha256()
    h.update(PROFILE_VERSION.encode())
    if pptx_path and Path(pptx_path).exists():
        h.update(b"PPTX:")
        h.update(Path(pptx_path).name.encode("utf-8"))
        with open(pptx_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
    for img in images:
        h.update(b"|IMG:")
        h.update(Path(img).name.encode("utf-8"))
        with open(img, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
    return h.hexdigest()[:16]


# =============================================================================
# VisionClient
# =============================================================================

class VisionClient:
    """OpenAI 兼容多模态 chat completions 客户端，独立读 VISION_* 三个环境变量。"""

    def __init__(self) -> None:
        self.base_url = os.getenv("VISION_BASE_URL", "").rstrip("/")
        self.api_key = os.getenv("VISION_API_KEY", "")
        self.model = os.getenv("VISION_MODEL_NAME", "gemini-3.1-pro-preview")
        if not self.base_url:
            raise ValueError("缺少 VISION_BASE_URL（请通过 agent 配置 / 系统环境变量注入）")
        if not self.api_key:
            raise ValueError("缺少 VISION_API_KEY（请通过 agent 配置 / 系统环境变量注入）")
        # 端点 URL：base 已包含 /v1 时不再追加
        if "/chat/completions" in self.base_url:
            self.endpoint = self.base_url
        elif self.base_url.rstrip("/").endswith("/v1"):
            self.endpoint = f"{self.base_url}/chat/completions"
        else:
            self.endpoint = f"{self.base_url}/v1/chat/completions"

    def chat_json(
        self,
        system: str,
        user_text: str,
        images: Optional[List[str]] = None,
        temperature: float = 0.2,
    ) -> Any:
        """发起 chat completions（非流式），强制要求模型返回 JSON。返回 parsed JSON。"""
        content: List[Dict[str, Any]] = []
        for img in images or []:
            content.append({"type": "image_url", "image_url": {"url": _data_url_for(img)}})
        content.append({"type": "text", "text": user_text})

        payload = {
            "model": self.model,
            "messages": [
                {"role": "system", "content": system},
                {"role": "user", "content": content},
            ],
            "temperature": temperature,
            "stream": False,
        }
        # response_format 在多数 OpenAI 兼容中转上是可选的，加上不会害；不支持时无视
        payload["response_format"] = {"type": "json_object"}

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
        }

        import time as _time
        last_err: Optional[Exception] = None
        for attempt in range(1, VISION_MAX_RETRIES + 1):
            try:
                resp = requests.post(
                    self.endpoint, headers=headers, json=payload, timeout=VISION_TIMEOUT_SECS
                )
                if resp.status_code != 200:
                    raise RuntimeError(f"vision 调用失败 status={resp.status_code} body={resp.text[:300]}")
                data = resp.json()
                msg = (
                    data.get("choices", [{}])[0]
                    .get("message", {})
                    .get("content", "")
                )
                if not msg:
                    raise RuntimeError(f"vision 返回空 content: {str(data)[:300]}")
                return _parse_json_loose(msg)
            except Exception as e:
                last_err = e
                msg_s = str(e)[:200]
                transient = any(s in msg_s for s in ("timeout", "Read timed out", "Connection aborted",
                                                     "RemoteDisconnected", "502", "503", "504", "524"))
                if attempt < VISION_MAX_RETRIES and transient:
                    print(f"(!)  vision 第 {attempt} 次失败({msg_s})，{VISION_RETRY_DELAY}s 后重试")
                    _time.sleep(VISION_RETRY_DELAY)
                    continue
                raise
        raise RuntimeError(f"vision 重试 {VISION_MAX_RETRIES} 次仍失败: {last_err}")


# =============================================================================
# Prompt 模板（中文，引导 Gemini 输出严格 JSON）
# =============================================================================

GLOBAL_SYSTEM = (
    "你是一名专业 PPT 视觉设计分析师。"
    "用户会给你一份 PPT 模板的若干页缩略图，你需要：\n"
    "1) 提炼整套模板的全局视觉风格（global_style）；\n"
    "2) 对每一页给出 page_type、布局摘要（summary）；\n"
    "3) 对每一页根据其内容容量给出严格的 JSON Schema（properties / required / minLength / maxLength / description）。\n"
    "**只返回 JSON，不要 markdown，不要任何额外解释**。"
)

GLOBAL_USER_TPL = """请基于下方 {n_images} 张 PPT 页的缩略图（按从前到后的顺序排列），输出 JSON：

{{
  "global_style": "...用 80-200 字概括整套模板的视觉调性，覆盖：背景/主色/强调色、字体观感（粗体/衬线/无衬线）、装饰元素、留白与构图风格...",
  "theme": {{
    "primary": "#hex",
    "accent": "#hex",
    "background": "#hex",
    "fonts": {{ "title": "字体名或 unknown", "body": "字体名或 unknown" }}
  }},
  "layouts": [
    {{
      "id": "唯一短 ID（kebab-case，描述布局，比如 cover-large-title / data-three-metrics / agenda-numbered-list）",
      "page_index": 0,
      "page_type": "cover|agenda|section|content|data|quote|closing|other",
      "summary": "用 60-120 字描述这页的视觉布局和内容编排方式，重点描写位置、装饰、文字层级，不要描述具体的文字内容",
      "visual_signature": "用 10-30 字概括这页最容易被观众记住的版式特征，如右图左文/三卡片/编号时间线",
      "content_capacity": {{
        "title": "建议标题字数或形态",
        "body": "建议正文容量，例如 3-5 个短要点 / 1 个大数字 + 2 条解释"
      }},
      "best_for": ["适合承载的内容结构，例如 agenda / feature list / comparison / quote"],
      "avoid_for": ["不适合承载的内容结构，例如 dense table / long legal copy"],
      "variation_tags": ["split", "cards", "timeline", "image-right"],
      "external_image_slots": [
        {
          "id": "hero-image",
          "purpose": "这块区域在模板里适合替换为真实图片/截图/图表",
          "bbox": [0.52, 0.14, 0.40, 0.68],
          "priority": 1
        }
      ],
      "reuse_friendly": true,
      "reuse_reason": "用 30-60 字解释为什么这页适合 / 不适合在同一份 deck 里被多次复用",
      "json_schema": {{
        "type": "object",
        "properties": {{
          "title":    {{ "type": "string", "minLength": 2, "maxLength": 24, "description": "..." }},
          "subtitle": {{ "type": "string", "minLength": 4, "maxLength": 50, "description": "..." }}
        }},
        "required": ["title"],
        "additionalProperties": false
      }}
    }}
  ]
}}

要求：
- layouts 数组的长度必须等于 {n_images}，每页一个，按图片顺序。
- 每个 layout.json_schema 必须严格按 JSON Schema 规范，properties 字段名用英文 key 但 description 用中文。
- 如果模板页里有明显照片区、图片区、图表区、截图区、主视觉空区、卡片图片区，请在 external_image_slots 中给出 1-3 个候选区域。
- external_image_slots[].bbox 使用归一化坐标 [x, y, w, h]，左上角为 0,0，右下角为 1,1；只表示适合放真实图的区域，不要写文字区。
- 如果模板页没有适合放真实图片的区域，external_image_slots 返回空数组 []。
- 字段长度 (minLength/maxLength) 要根据该页可容纳的字数实际给出，不要给死值。
- 如果某页含数据图表，必须有 "metrics" 这种 array 字段；如果是列表页要有 "items" array；
  封面要有 title/subtitle；目录页要有 items；数据页要有 metrics 或 stats。
- visual_signature / content_capacity / best_for / avoid_for / variation_tags 用来帮助后续自动分配不同页面形态：
  * visual_signature 要短，描述视觉差异，而不是复述 page_type。
  * content_capacity 要表达这页最多适合放多少信息，避免把长文塞进强视觉页。
  * variation_tags 使用英文短标签，方便比较布局差异。
- 颜色用 #RRGGBB；不确定就给最接近的近似值。
- 严格只返回 JSON 一个对象，无其他文本。

【reuse_friendly 判定规则】
- true（可重复使用）：版式通用、装饰元素与具体语义无强绑定。典型例子：
  * 纯文字 / 多卡片网格 / 多条目列表 / 通用数据对比 / 章节小节标题
  * 装饰只是几何形状 / 抽象色块 / 通用图标
- false（不建议重复）：版式带强语义锚点，复用会让观众觉得"为什么又是这页"。典型例子：
  * 封面页（cover）-- 整份 deck 只能有 1 个
  * 含独特角色插画的页（如 3 个具名人物 + 不同职业）
  * 含独特场景插画的页（如雪山、广播塔、复古收音机等显眼意象）
  * 多步骤流程页 + 每步独有图标（如 5 步骤 zigzag with 不同 icon），复用会暗示"同一个流程"
  * novelty 数据可视化页（独特中央装置图）
- reuse_reason 用一句话讲明白判定依据，方便用户决策。
"""


def vision_analyze(images: List[str], client: VisionClient, pptx_meta: Dict[str, Any]) -> Dict[str, Any]:
    """一次性把所有页缩略图传给 vision，得到完整 TemplateProfile（不含 source/cache 字段）。"""
    if not images:
        raise ValueError("vision_analyze 需要至少一张模板图片")
    user_text = GLOBAL_USER_TPL.format(n_images=len(images))
    if pptx_meta:
        user_text += f"\n\n参考：模板原 .pptx 比例 = {pptx_meta.get('aspect', 'unknown')}，共 {pptx_meta.get('page_count', len(images))} 页。"

    print(f"🔍 vision 分析 {len(images)} 页模板（model={client.model}）...")
    result = client.chat_json(
        system=GLOBAL_SYSTEM,
        user_text=user_text,
        images=images,
        temperature=0.2,
    )
    if not isinstance(result, dict) or "layouts" not in result:
        raise ValueError(f"vision 返回结构不合法: {str(result)[:300]}")

    layouts = result.get("layouts", [])
    if len(layouts) > len(images):
        print(f"(!)  vision 返回 {len(layouts)} 个 layout，超过图片数 {len(images)}；截断")
        layouts = layouts[: len(images)]
    if len(layouts) < len(images):
        # vision 经常因 token 限制少返回。补齐占位 layout，把对应模板图当 reference image。
        missing = len(images) - len(layouts)
        print(f"(!)  vision 返回 {len(layouts)} 个 layout，少于图片数 {len(images)}；补齐 {missing} 个占位")
        for i in range(len(layouts), len(images)):
            layouts.append({
                "id": f"layout-{i + 1:02d}",
                "page_index": i,
                "page_type": "content",
                "summary": "（vision 未返回此页 layout，直接用模板图作为参考）",
                "reuse_friendly": True,
                "reuse_reason": "（vision 未分析，默认按可复用处理）",
                "json_schema": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string", "minLength": 2, "maxLength": 40},
                        "body": {"type": "string", "minLength": 0, "maxLength": 600},
                    },
                    "required": ["title"],
                    "additionalProperties": False,
                },
            })
    result["layouts"] = layouts

    # 给每个 layout 注入 reference_image 路径并校准基础字段
    for i, layout in enumerate(layouts):
        layout["reference_image"] = str(Path(images[i]).resolve())
        layout.setdefault("page_index", i)
        pt = layout.get("page_type", "")
        if pt not in VALID_PAGE_TYPES:
            layout["page_type"] = "content"
        layout.setdefault("id", f"layout-{i + 1:02d}")
        # 复用友好性：cover 永远视为不可复用；其他默认 True（vision 没返回时不打扰用户）
        if layout.get("page_type") == "cover":
            layout["reuse_friendly"] = False
            layout.setdefault("reuse_reason", "封面页一份 deck 只能有 1 个")
        else:
            layout.setdefault("reuse_friendly", True)
            layout.setdefault("reuse_reason", "")
        layout["visual_signature"] = str(layout.get("visual_signature") or "").strip()
        capacity = layout.get("content_capacity")
        layout["content_capacity"] = capacity if isinstance(capacity, (dict, list, str)) else {}
        for key in ("best_for", "avoid_for", "variation_tags"):
            val = layout.get(key)
            if isinstance(val, list):
                layout[key] = [str(x).strip() for x in val if str(x).strip()]
            elif isinstance(val, str) and val.strip():
                layout[key] = [val.strip()]
            else:
                layout[key] = []
        slots = layout.get("external_image_slots")
        layout["external_image_slots"] = slots if isinstance(slots, list) else []
    return result


# =============================================================================
# 顶层入口（含缓存）
# =============================================================================

def _cache_dir() -> Path:
    """模板 vision 缓存放在调用者 cwd 下，不污染 skill 安装目录。"""
    return Path.cwd() / CACHE_DIR_NAME


def analyze_template(
    pptx_path: Optional[str],
    images_dir: Optional[str],
    rebuild: bool = False,
    client: Optional[VisionClient] = None,
    cache_dir: Optional[Path] = None,
) -> Dict[str, Any]:
    """主入口：返回完整 TemplateProfile。

    - 仅 pptx_path、未传图片：fallback，只返回基础 metadata，不含 layouts（无法做仿模板）。
    - 同时传图片 -> 调 vision 出完整 profile，命中缓存直接读。
    """
    pptx_meta = read_pptx_metadata(pptx_path) if pptx_path else {}
    images = load_template_images(images_dir) if images_dir else []

    source_label = Path(pptx_path).name if pptx_path else (Path(images_dir).name if images_dir else "unknown")

    if not images:
        print("(!)  未传 --template-images，无法做 vision 风格分析，仅返回 .pptx 元数据")
        return {
            "version": PROFILE_VERSION,
            "source": source_label,
            "source_hash": compute_source_hash(pptx_path, []),
            "global_style": "",
            "theme": {},
            "layouts": [],
            "pptx_meta": pptx_meta,
        }

    src_hash = compute_source_hash(pptx_path, images)
    cache_root = cache_dir if cache_dir else _cache_dir()
    cache_root.mkdir(parents=True, exist_ok=True)
    cache_path = cache_root / f"{src_hash}.json"

    if cache_path.exists() and not rebuild:
        print(f"📦 模板缓存命中 {cache_path}")
        try:
            with open(cache_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            print(f"(!)  缓存损坏 ({e})，重新分析")

    if client is None:
        client = VisionClient()
    profile = vision_analyze(images, client, pptx_meta)

    profile["version"] = PROFILE_VERSION
    profile["source"] = source_label
    profile["source_hash"] = src_hash
    profile["pptx_meta"] = pptx_meta

    # 防御性 mkdir：极少数情况下 cache_root 在前面 mkdir 后被外部清掉
    cache_path.parent.mkdir(parents=True, exist_ok=True)
    with open(cache_path, "w", encoding="utf-8") as f:
        json.dump(profile, f, ensure_ascii=False, indent=2)
    print(f"💾 模板 profile 已缓存 -> {cache_path}")
    return profile


# =============================================================================
# layout 匹配 + 字段拆解
# =============================================================================

def match_layout(slide: Dict[str, Any], profile: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """按 layout_id -> layout-NN 数字尾巴 -> page_type -> slide_number 索引兜底，返回最贴合 slide 的 layout dict。"""
    layouts = profile.get("layouts", [])
    if not layouts:
        return None

    layout_id = slide.get("layout_id")
    if layout_id:
        for lay in layouts:
            if lay.get("id") == layout_id:
                return lay
        # layout_id 找不到 -> 如果是 layout-NN 这种自动 ID，按数字尾巴兜底到第 NN 个 layout，
        # 保留用户"slide N -> 模板第 N 页"的意图
        m = re.search(r"(\d+)\s*$", layout_id)
        if m:
            idx = int(m.group(1)) - 1
            if 0 <= idx < len(layouts):
                print(f"(!)  slide {slide.get('slide_number')} 的 layout_id={layout_id} 缺失，按编号回退 layouts[{idx}]")
                return layouts[idx]
        print(f"(!)  slide {slide.get('slide_number')} 指定的 layout_id={layout_id} 在模板中不存在，回退 page_type")

    page_type = slide.get("page_type", "content")
    type_matches = [lay for lay in layouts if lay.get("page_type") == page_type]
    if type_matches:
        return type_matches[0]

    # 都找不到 -> 用第几页对位（slide_number 对应模板第几页，超出取最后）
    n = slide.get("slide_number", 1)
    idx = max(0, min(n - 1, len(layouts) - 1))
    return layouts[idx]


def _routing_content_text(slide: Dict[str, Any]) -> str:
    content = slide.get("content", "")
    if isinstance(content, str):
        return content
    try:
        return json.dumps(content, ensure_ascii=False)
    except (TypeError, ValueError):
        return str(content)


def _routing_features(slide: Dict[str, Any]) -> Dict[str, Any]:
    """Extract conservative content-shape hints for opt-in distilled layouts."""
    text = _routing_content_text(slide)
    lowered = text.lower()
    metric_count = len(re.findall(r"[-+]?\d+(?:\.\d+)?\s*%", text))
    quarter_count = len(set(re.findall(r"\bq[1-4]\b", lowered)))
    year_count = len(set(re.findall(r"\b(?:19|20)\d{2}\b", text)))
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    bullet_count = sum(
        1 for line in lines if re.match(r"^(?:[-*•]|\d+[.)、])\s*", line)
    )
    item_count = max(metric_count, bullet_count, max(0, len(lines) - 1))
    shapes: set[str] = set()
    if metric_count:
        shapes.add("metrics")
    if metric_count >= 2 and quarter_count >= 2:
        shapes.update({"metrics-series", "metrics+series"})
    if year_count >= 2 or any(token in lowered for token in ("timeline", "时间线", "里程碑")):
        shapes.add("timeline")
    if any(token in lowered for token in ("before", "after", "versus", " vs ", "对比", "之前", "之后")):
        shapes.add("comparison")
    if any(token in lowered for token in ("step", "phase", "阶段", "步骤", "流程")):
        shapes.add("process")
    if "|" in text or any(token in lowered for token in ("table", "表格", "矩阵")):
        shapes.add("table")
    if bullet_count >= 2 or len(lines) >= 4:
        shapes.add("list")
    if not shapes:
        shapes.add("general")
    return {
        "text": lowered,
        "shapes": shapes,
        "metric_count": metric_count,
        "series_count": max(quarter_count, year_count),
        "item_count": item_count,
    }


def _routing_shape_aliases(values: set[str]) -> set[str]:
    """Normalize model-authored shape synonyms into stable routing families."""
    expanded: set[str] = set()
    for value in values:
        token = re.sub(r"[^a-z0-9]+", "-", str(value).strip().lower()).strip("-")
        if not token:
            continue
        expanded.add(token)
        if any(part in token for part in ("metric", "kpi", "chart", "series", "trend")):
            expanded.add("metrics")
        if any(part in token for part in ("table", "tabular", "matrix")):
            expanded.add("table")
        if any(part in token for part in ("timeline", "milestone", "roadmap", "chronolog")):
            expanded.add("timeline")
        if any(part in token for part in ("compar", "versus")):
            expanded.add("comparison")
        if any(part in token for part in ("process", "workflow", "step", "phase")):
            expanded.add("process")
        if any(part in token for part in ("list", "bullet")):
            expanded.add("list")
    return expanded


def layout_routing_score(slide: Dict[str, Any], layout: Dict[str, Any]) -> float:
    """Score explicit routing metadata; layouts without routing remain order-driven."""
    routing = layout.get("routing")
    if not isinstance(routing, dict) or not routing:
        return 0.0
    features = _routing_features(slide)
    shapes = _routing_shape_aliases(set(features["shapes"]))
    score = 1.0
    content_shapes = routing.get("content_shapes")
    if isinstance(content_shapes, str):
        content_shapes = [content_shapes]
    if isinstance(content_shapes, list) and content_shapes:
        wanted = _routing_shape_aliases(
            {str(item).strip().lower() for item in content_shapes if str(item).strip()}
        )
        score += 30.0 if shapes & wanted else -15.0
    for field, delta_present, delta_missing in (
        ("requires", 10.0, -25.0),
        ("excludes", -30.0, 0.0),
    ):
        values = routing.get(field)
        if isinstance(values, str):
            values = [values]
        if not isinstance(values, list):
            continue
        for value in values:
            token = str(value).strip().lower()
            if not token:
                continue
            present = bool(_routing_shape_aliases({token}) & shapes) or token in features["text"]
            score += delta_present if present else delta_missing
    for feature_name, min_key, max_key in (
        ("item_count", "min_items", "max_items"),
        ("metric_count", "min_metrics", "max_metrics"),
        ("series_count", "min_series", "max_series"),
    ):
        value = int(features[feature_name])
        minimum = routing.get(min_key)
        maximum = routing.get(max_key)
        if isinstance(minimum, (int, float)):
            score += 4.0 if value >= minimum else -12.0
        if isinstance(maximum, (int, float)):
            score += 4.0 if value <= maximum else -12.0
    for key, delta in (("keywords", 3.0), ("exclude_keywords", -20.0)):
        values = routing.get(key)
        if isinstance(values, str):
            values = [values]
        if isinstance(values, list):
            for value in values:
                token = str(value).strip().lower()
                if token and token in features["text"]:
                    score += delta
    return score


def assign_layouts(slides: List[Dict[str, Any]], profile: Dict[str, Any]) -> Dict[int, Dict[str, Any]]:
    """Assign template layouts to slides while preferring one layout per page.

    Explicit `layout_id` always wins, including intentional reuse. Slides without
    an explicit layout get the best unused layout for their page_type before
    falling back to the older match_layout behavior.
    """
    layouts = profile.get("layouts", [])
    if not layouts:
        return {}

    assigned: Dict[int, Dict[str, Any]] = {}
    used_ids = set()

    def layout_id_of(layout: Dict[str, Any], idx: int) -> str:
        return str(layout.get("id") or f"layout-{idx + 1:02d}")

    def mark(slide: Dict[str, Any], layout: Dict[str, Any]) -> None:
        slide_number = int(slide.get("slide_number", len(assigned) + 1))
        assigned[slide_number] = layout
        try:
            idx = layouts.index(layout)
        except ValueError:
            idx = 0
        used_ids.add(layout_id_of(layout, idx))

    # Pass 1: honor explicit layout choices exactly.
    for slide in slides:
        layout_id = slide.get("layout_id")
        if not layout_id:
            continue
        chosen = None
        for lay in layouts:
            if lay.get("id") == layout_id:
                chosen = lay
                break
        if chosen is None:
            m = re.search(r"(\d+)\s*$", str(layout_id))
            if m:
                idx = int(m.group(1)) - 1
                if 0 <= idx < len(layouts):
                    print(f"(!)  slide {slide.get('slide_number')} 的 layout_id={layout_id} 缺失，按编号回退 layouts[{idx}]")
                    chosen = layouts[idx]
        if chosen is None:
            print(f"(!)  slide {slide.get('slide_number')} 指定的 layout_id={layout_id} 在模板中不存在，自动分配未使用 layout")
            continue
        mark(slide, chosen)

    # Pass 2: assign unused layouts, preserving page_type where possible.
    for slide in slides:
        slide_number = int(slide.get("slide_number", len(assigned) + 1))
        if slide_number in assigned:
            continue

        page_type = slide.get("page_type", "content")
        chosen = None

        unused_same_type = [
            lay
            for i, lay in enumerate(layouts)
            if lay.get("page_type") == page_type and layout_id_of(lay, i) not in used_ids
        ]
        if any(isinstance(lay.get("routing"), dict) and lay.get("routing") for lay in unused_same_type):
            chosen = max(
                unused_same_type,
                key=lambda lay: layout_routing_score(slide, lay),
            )

        if chosen is None:
            idx = max(0, min(slide_number - 1, len(layouts) - 1))
            indexed = layouts[idx]
            indexed_id = layout_id_of(indexed, idx)
            if indexed.get("page_type") == page_type and indexed_id not in used_ids:
                chosen = indexed

        if chosen is None:
            for i, lay in enumerate(layouts):
                lid = layout_id_of(lay, i)
                if lay.get("page_type") == page_type and lid not in used_ids:
                    chosen = lay
                    break

        same_type = [lay for lay in layouts if lay.get("page_type") == page_type]
        if chosen is None and same_type:
            reusable = [lay for lay in same_type if lay.get("reuse_friendly", True)]
            if any(isinstance(lay.get("routing"), dict) and lay.get("routing") for lay in reusable):
                chosen = max(reusable, key=lambda lay: layout_routing_score(slide, lay))
            else:
                for lay in reusable:
                    chosen = lay
                    break
            if chosen is None:
                chosen = same_type[0]

        if chosen is None:
            for i, lay in enumerate(layouts):
                lid = layout_id_of(lay, i)
                if lid not in used_ids:
                    chosen = lay
                    break

        if chosen is None:
            # All layouts are already used; reuse the best old-style match.
            tmp = dict(slide)
            tmp.pop("layout_id", None)
            chosen = match_layout(tmp, profile)

        if chosen:
            mark(slide, chosen)

    return assigned


def check_layout_reuse(slides: List[Dict[str, Any]], profile: Dict[str, Any]) -> List[str]:
    """检测 plan 里同一 layout 被多次使用的情况，并指出未被使用的"备胎"。

    返回一个警告列表（已格式化的字符串），便于上层直接 print。
    判定路径：先用 match_layout 拿到每个 slide 实际命中的 layout，再按
    layout.id 聚合。
    - 重复 + reuse_friendly == False：(!) 强警告，建议必须换；
    - 重复 + reuse_friendly == True：(i) 弱提示，建议优先用没用过的 layout；
    - 同时附带还未被使用的候选 layout（按 page_type 归组），方便决策。
    """
    warnings: List[str] = []
    layouts = profile.get("layouts", [])
    if not layouts:
        return warnings

    # 收集每个 layout_id 被哪些 slide 用了
    assigned = assign_layouts(slides, profile)
    usage: Dict[str, List[int]] = {}
    layout_meta: Dict[str, Dict[str, Any]] = {}
    for slide in slides:
        lay = assigned.get(int(slide.get("slide_number", 0)))
        if not lay:
            continue
        lid = lay.get("id", "?")
        usage.setdefault(lid, []).append(slide.get("slide_number", 0))
        layout_meta[lid] = lay

    used_ids = set(usage.keys())
    unused = [lay for lay in layouts if lay.get("id") not in used_ids]

    for lid, slide_nums in usage.items():
        if len(slide_nums) < 2:
            continue
        lay = layout_meta[lid]
        friendly = lay.get("reuse_friendly", True)
        page_type = lay.get("page_type", "?")
        reason = lay.get("reuse_reason") or "建议尽量做到一页一种 layout"
        # 找同 page_type 还没用的 layout 做候选
        candidates = [u.get("id") for u in unused if u.get("page_type") == page_type][:3]
        candidate_hint = (
            f"\n    备胎 layout（同 page_type={page_type}，可考虑切到这些）：{', '.join(candidates)}"
            if candidates else ""
        )
        if friendly:
            warnings.append(
                f"(i)  layout {lid}（{page_type}）被 slide {slide_nums} 重复 {len(slide_nums)} 次使用，"
                f"虽标 reuse_friendly=True，仍建议换不同 layout 避免视觉重复。{candidate_hint}"
            )
        else:
            warnings.append(
                f"(!)  layout {lid}（{page_type}）被 slide {slide_nums} 重复 {len(slide_nums)} 次使用，"
                f"且标 reuse_friendly=False。原因：{reason}\n"
                f"    强烈建议换 layout。{candidate_hint}"
            )
    return warnings


def summarize_layouts(profile: Dict[str, Any]) -> str:
    """把 profile 里的所有 layout 整理成一段表格状文本，方便人/Claude 选页。"""
    layouts = profile.get("layouts", [])
    if not layouts:
        return "（profile 没有 layout）"
    lines = ["id  | page_type | reuse | summary"]
    lines.append("--- | --- | --- | ---")
    for lay in layouts:
        flag = "✓" if lay.get("reuse_friendly", True) else "✗"
        summary = (lay.get("summary", "") or "").replace("\n", " ")[:60]
        lines.append(f"{lay.get('id','?')} | {lay.get('page_type','?')} | {flag} | {summary}")
    return "\n".join(lines)


def coerce_fields(slide: Dict[str, Any], layout: Dict[str, Any]) -> Dict[str, Any]:
    """如果 slide 已含 fields，按 schema 验证；否则把 content 字符串拆成 fields。

    拆解策略：
    - 第一行作为 title
    - 其余非空行作为 items / bullets / body
    - 严格按 schema.properties 的字段名和类型
    """
    schema = layout.get("json_schema") or {"type": "object", "properties": {}}
    properties = (schema.get("properties") or {}) if isinstance(schema, dict) else {}

    if "fields" in slide and isinstance(slide["fields"], dict):
        fields = dict(slide["fields"])
    else:
        fields = _content_to_fields(slide.get("content", ""), properties)

    # 简单截断（避免超 maxLength 触发 schema validation 失败）
    for key, prop in properties.items():
        if key not in fields:
            continue
        v = fields[key]
        max_len = prop.get("maxLength") if isinstance(prop, dict) else None
        if isinstance(v, str) and isinstance(max_len, int) and len(v) > max_len:
            fields[key] = v[:max_len]

    # jsonschema 软校验
    try:
        import jsonschema
        try:
            jsonschema.validate(instance=fields, schema=schema)
        except jsonschema.ValidationError as ve:
            print(f"(!)  slide {slide.get('slide_number')} 字段未通过 schema: {ve.message[:120]}")
    except ImportError:
        pass
    return fields


def _content_to_fields(content: str, properties: Dict[str, Any]) -> Dict[str, Any]:
    """把自由 content 字符串按 schema properties 切分到字段。"""
    if not content:
        return {}
    lines = [ln.strip() for ln in content.splitlines() if ln.strip()]
    if not lines:
        return {}

    fields: Dict[str, Any] = {}
    title = lines[0]
    rest_lines = lines[1:]

    # 标准字段：title / subtitle
    if "title" in properties:
        fields["title"] = title
    if "subtitle" in properties and rest_lines:
        fields["subtitle"] = rest_lines[0]
        rest_lines = rest_lines[1:]
    if "tagline" in properties and rest_lines:
        fields["tagline"] = rest_lines[0]
        rest_lines = rest_lines[1:]

    # array 字段：items / metrics / bullets / points
    array_field = None
    array_item_props: Dict[str, Any] = {}
    for k, p in properties.items():
        if isinstance(p, dict) and p.get("type") == "array":
            array_field = k
            items_schema = p.get("items") or {}
            if isinstance(items_schema, dict):
                array_item_props = items_schema.get("properties") or {}
            break

    if array_field and rest_lines:
        items: List[Any] = []
        for ln in rest_lines:
            cleaned = re.sub(r"^[\-\*\u2022*•\d\.\)\s]+", "", ln).strip()
            if not cleaned:
                continue
            if not array_item_props:
                items.append(cleaned)
                continue
            obj: Dict[str, Any] = {}
            # 拆 "name: desc" / "name - desc" / "name | desc"
            split_match = re.split(r"[：:\-\|]", cleaned, maxsplit=1)
            name_val = split_match[0].strip() if split_match else cleaned
            desc_val = split_match[1].strip() if len(split_match) > 1 else ""
            for prop_key in array_item_props.keys():
                if prop_key in ("name", "label", "title", "key"):
                    obj[prop_key] = name_val
                elif prop_key in ("value", "metric", "number"):
                    # 抽数字+单位
                    m = re.search(r"([\d\.,%]+\s*[万千亿%倍kKmMbB+]*)", cleaned)
                    obj[prop_key] = m.group(1) if m else cleaned
                else:
                    obj[prop_key] = desc_val or cleaned
            items.append(obj)
        if items:
            fields[array_field] = items
    elif rest_lines:
        # 没 array 字段：把剩余文本拼成 body / description
        for k in ("body", "description", "content", "text"):
            if k in properties:
                fields[k] = "\n".join(rest_lines)
                break

    return fields


# =============================================================================
# 字段值 -> 图片 prompt 拼接
# =============================================================================

def render_prompt_from_template(
    profile: Dict[str, Any],
    layout: Dict[str, Any],
    fields: Dict[str, Any],
    aspect_hint: str = "",
    language_rule: str = "",
) -> str:
    """按 plan 里的拼接规则组装最终图片 prompt。"""
    parts: List[str] = []
    g_style = profile.get("global_style", "").strip()
    if g_style:
        parts.append(f"【全局风格】\n{g_style}")
    parts.append(f"【布局描述】\n{layout.get('summary', '').strip()}")

    variation_lines: List[str] = []
    visual_signature = str(layout.get("visual_signature") or "").strip()
    if visual_signature:
        variation_lines.append(f"- 视觉签名: {visual_signature}")
    content_capacity = layout.get("content_capacity")
    if content_capacity:
        if isinstance(content_capacity, str):
            capacity_text = content_capacity
        else:
            capacity_text = json.dumps(content_capacity, ensure_ascii=False)
        variation_lines.append(f"- 内容容量: {capacity_text}")
    for key, label in (
        ("best_for", "适合内容"),
        ("avoid_for", "避免内容"),
        ("variation_tags", "变体标签"),
    ):
        val = layout.get(key)
        if isinstance(val, list) and val:
            variation_lines.append(f"- {label}: {', '.join(str(x) for x in val if str(x))}")
        elif isinstance(val, str) and val.strip():
            variation_lines.append(f"- {label}: {val.strip()}")
    if variation_lines:
        parts.append("【版式变体信息】\n" + "\n".join(variation_lines))

    parts.append("【本页内容】")
    schema_props = ((layout.get("json_schema") or {}).get("properties") or {}) if isinstance(layout.get("json_schema"), dict) else {}
    for key, val in fields.items():
        prop = schema_props.get(key, {}) if isinstance(schema_props, dict) else {}
        desc = prop.get("description", "") if isinstance(prop, dict) else ""
        if isinstance(val, list):
            rendered_items = []
            for item in val:
                if isinstance(item, dict):
                    rendered_items.append(" / ".join(f"{k}={v}" for k, v in item.items() if v))
                else:
                    rendered_items.append(str(item))
            joined = "\n  - " + "\n  - ".join(rendered_items)
            parts.append(f"- {key}（{desc}）：{joined}")
        else:
            parts.append(f"- {key}（{desc}）：{val}")

    if language_rule:
        parts.append(language_rule)
    if aspect_hint:
        parts.append(aspect_hint)

    return "\n\n".join(parts)


# =============================================================================
# CLI 自检（可选）
# =============================================================================

if __name__ == "__main__":
    import argparse
    from dotenv import load_dotenv

    parser = argparse.ArgumentParser(description="Template Analyzer CLI (debug)")
    parser.add_argument("--pptx", help="模板 .pptx 路径（可选）")
    parser.add_argument("--images", required=True, help="模板每页 PNG 所在目录")
    parser.add_argument("--rebuild", action="store_true", help="无视缓存重新分析")
    args = parser.parse_args()

    # 找 .env
    for cand in [
        Path(__file__).parent / ".env",
        Path.home() / ".claude/skills/gpt-image2-ppt-skills/.env",
    ]:
        if cand.exists():
            load_dotenv(cand, override=True)
            print(f"loaded env: {cand}")
            break

    profile = analyze_template(args.pptx, args.images, rebuild=args.rebuild)
    print(json.dumps(profile, ensure_ascii=False, indent=2)[:3000])
